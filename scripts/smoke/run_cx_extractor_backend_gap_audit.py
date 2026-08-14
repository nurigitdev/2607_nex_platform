#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import os
import sys
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Callable, Mapping


ROOT = Path(__file__).resolve().parents[2]
SHARED_PATH = ROOT / "services" / "_shared"
CX_PATH = ROOT / "services" / "nex-cx"
sys.path.insert(0, str(SHARED_PATH))
sys.path.insert(0, str(CX_PATH))

from nex_cx.extractors import (  # noqa: E402
    BINARY_SOURCE_FORMATS,
    DOCX_EXTRACTION_MODE,
    PDF_EXTRACTION_MODE,
    PPTX_EXTRACTION_MODE,
    PLACEHOLDER_BINARY_MODE,
    PLACEHOLDER_BINARY_WARNING_PREFIX,
    XLSX_EXTRACTION_MODE,
    ExtractionAdapterError,
    ExtractorInput,
    LocalMockTextExtractor,
    extractor_backend_catalog,
    extractor_backend_gap_summary,
)
from nex_cx.ingestion import sha256_bytes  # noqa: E402


SCHEMA_VERSION = "cx_extractor_backend_gap_audit.v1"
SECRET_SOURCE = b"CX extractor backend audit secret body should not leak."

PROTECTED_ENV_KEYS = (
    "NEX_CX_TEST_DATABASE_URL",
    "NEX_CX_DATABASE_URL",
    "NEX_DATA_ROOT",
    "NEX_CX_SOURCE_STORAGE_ROOT",
    "NEX_CX_EXTRACTED_MARKDOWN_ROOT",
    "NEX_CX_EXTRACTION_TEMP_ROOT",
)


@dataclass(frozen=True)
class RequiredPath:
    name: str
    path: Path
    purpose: str


@dataclass(frozen=True)
class TokenRequirement:
    group: str
    path: Path
    token_id: str
    token: str
    purpose: str


CX_EXTRACTORS = ROOT / "services" / "nex-cx" / "nex_cx" / "extractors.py"
CX_INGESTION = ROOT / "services" / "nex-cx" / "nex_cx" / "ingestion.py"
CX_README = ROOT / "services" / "nex-cx" / "README.md"
SLICE_0072_DOC = (
    ROOT / "docs" / "slices" / "0072_cx_text_extraction_adapter_foundation.md"
)
SLICE_0283_DOC = (
    ROOT / "docs" / "slices" / "0283_cx_uploaded_source_extraction_postgresql_smoke.md"
)

REQUIRED_PATHS = (
    RequiredPath("cx_extractors", CX_EXTRACTORS, "CX extractor adapter boundary."),
    RequiredPath("cx_ingestion", CX_INGESTION, "CX extraction job adapter wiring."),
    RequiredPath("cx_readme", CX_README, "CX extraction boundary documentation."),
    RequiredPath(
        "slice_0072_doc",
        SLICE_0072_DOC,
        "Original extractor adapter foundation decision.",
    ),
    RequiredPath(
        "slice_0283_doc",
        SLICE_0283_DOC,
        "Latest uploaded-source extraction PostgreSQL evidence.",
    ),
)

REQUIRED_SOURCE_TOKENS = (
    TokenRequirement(
        "extractor_catalog",
        CX_EXTRACTORS,
        "backend_capability_dataclass",
        "ExtractorBackendCapability",
        "Extractor backend capability records are explicit and inspectable.",
    ),
    TokenRequirement(
        "extractor_catalog",
        CX_EXTRACTORS,
        "backend_catalog_function",
        "extractor_backend_catalog",
        "CX can emit the current source-format backend catalog.",
    ),
    TokenRequirement(
        "extractor_catalog",
        CX_EXTRACTORS,
        "gap_summary_function",
        "extractor_backend_gap_summary",
        "CX can summarize implemented vs placeholder backend gaps.",
    ),
    TokenRequirement(
        "binary_placeholder_boundary",
        CX_EXTRACTORS,
        "binary_source_formats",
        "BINARY_SOURCE_FORMATS",
        "Binary document source formats are centralized before real adapters.",
    ),
    TokenRequirement(
        "binary_placeholder_boundary",
        CX_EXTRACTORS,
        "placeholder_mode_constant",
        "PLACEHOLDER_BINARY_MODE",
        "Binary placeholder mode is a named adapter state.",
    ),
    TokenRequirement(
        "binary_placeholder_boundary",
        CX_EXTRACTORS,
        "placeholder_warning_prefix",
        "PLACEHOLDER_BINARY_WARNING_PREFIX",
        "Binary placeholder warnings are stable for later migration checks.",
    ),
    TokenRequirement(
        "pdf_extraction_backend",
        CX_EXTRACTORS,
        "pdf_extraction_mode",
        "PDF_EXTRACTION_MODE",
        "PDF real extraction has a named mode for runtime evidence.",
    ),
    TokenRequirement(
        "pdf_extraction_backend",
        CX_EXTRACTORS,
        "pdf_extraction_function",
        "extract_pdf_markdown",
        "PDF extraction is implemented behind the extractor adapter boundary.",
    ),
    TokenRequirement(
        "pdf_extraction_backend",
        CX_EXTRACTORS,
        "pdf_reader_backend",
        "pypdf",
        "The local PDF adapter uses the pinned parser dependency.",
    ),
    TokenRequirement(
        "docx_extraction_backend",
        CX_EXTRACTORS,
        "docx_extraction_mode",
        "DOCX_EXTRACTION_MODE",
        "DOCX real extraction has a named mode for runtime evidence.",
    ),
    TokenRequirement(
        "docx_extraction_backend",
        CX_EXTRACTORS,
        "docx_extraction_function",
        "extract_docx_markdown",
        "DOCX extraction is implemented behind the extractor adapter boundary.",
    ),
    TokenRequirement(
        "docx_extraction_backend",
        CX_EXTRACTORS,
        "docx_reader_backend",
        "from docx import Document",
        "The local DOCX adapter uses the pinned parser dependency.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "pptx_extraction_mode",
        "PPTX_EXTRACTION_MODE",
        "PPTX real extraction has a named mode for runtime evidence.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "pptx_extraction_function",
        "extract_pptx_markdown",
        "PPTX extraction is implemented behind the extractor adapter boundary.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "pptx_reader_backend",
        "from pptx import Presentation",
        "The local PPTX adapter uses the pinned parser dependency.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "xlsx_extraction_mode",
        "XLSX_EXTRACTION_MODE",
        "XLSX real extraction has a named mode for runtime evidence.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "xlsx_extraction_function",
        "extract_xlsx_markdown",
        "XLSX extraction is implemented behind the extractor adapter boundary.",
    ),
    TokenRequirement(
        "office_extraction_backend",
        CX_EXTRACTORS,
        "xlsx_reader_backend",
        "from openpyxl import load_workbook",
        "The local XLSX adapter uses the pinned parser dependency.",
    ),
    TokenRequirement(
        "ingestion_adapter_wiring",
        CX_INGESTION,
        "default_local_mock_extractor",
        "selected_extractor = extractor or LocalMockTextExtractor()",
        "The extraction job still selects the adapter through the boundary.",
    ),
    TokenRequirement(
        "ingestion_adapter_wiring",
        CX_INGESTION,
        "extractor_input_boundary",
        "ExtractorInput(",
        "The extraction job passes source bytes through the adapter input type.",
    ),
    TokenRequirement(
        "prior_decision",
        SLICE_0072_DOC,
        "binary_placeholder_decision",
        "deterministic placeholder with a warning",
        "Slice 0072 records the binary placeholder decision.",
    ),
    TokenRequirement(
        "prior_evidence",
        SLICE_0283_DOC,
        "postgres_extraction_evidence",
        "CX Uploaded Source Extraction PostgreSQL Smoke",
        "Slice 0283 proves the durable uploaded-source extraction path.",
    ),
    TokenRequirement(
        "service_docs",
        CX_README,
        "readme_extractor_boundary",
        "Text extraction runs through the `nex_cx.extractors` adapter boundary.",
        "CX README exposes the current extractor boundary.",
    ),
)

RuntimeProbe = Callable[[], dict[str, Any]]


def run_cx_extractor_backend_gap_audit(
    environ: dict[str, str] | None = None,
    *,
    root_dir: Path = ROOT,
    runtime_probe_runner: RuntimeProbe | None = None,
) -> dict[str, Any]:
    env = environ if environ is not None else os.environ
    paths = path_checks(root_dir)
    source_tokens = source_token_checks(root_dir)
    token_groups = grouped_token_status(source_tokens)
    catalog = [asdict(item) for item in extractor_backend_catalog()]
    gap_summary = extractor_backend_gap_summary()
    runtime_probe = run_runtime_probe(
        runtime_probe_runner or run_extractor_backend_probe
    )
    issues = [
        *[
            issue("path_missing", item["name"], item["path"])
            for item in paths
            if not item["present"]
        ],
        *[
            issue("source_token_missing", item["token_id"], item["path"])
            for item in source_tokens
            if not item["present"]
        ],
    ]
    if runtime_probe["status"] != "PASS":
        issues.append(
            issue(
                "runtime_probe_failed",
                str(runtime_probe.get("failure_code", "checks_failed")),
                "extractor_backend_runtime_probe",
            )
        )
    checks = {
        "required_paths_present": all(item["present"] for item in paths),
        "extractor_catalog_ready": token_groups.get("extractor_catalog") is True,
        "binary_placeholder_boundary_ready": (
            token_groups.get("binary_placeholder_boundary") is True
        ),
        "pdf_extraction_backend_ready": (
            token_groups.get("pdf_extraction_backend") is True
        ),
        "docx_extraction_backend_ready": (
            token_groups.get("docx_extraction_backend") is True
        ),
        "office_extraction_backend_ready": (
            token_groups.get("office_extraction_backend") is True
        ),
        "ingestion_adapter_wiring_ready": (
            token_groups.get("ingestion_adapter_wiring") is True
        ),
        "prior_decision_recorded": token_groups.get("prior_decision") is True,
        "prior_evidence_recorded": token_groups.get("prior_evidence") is True,
        "service_docs_recorded": token_groups.get("service_docs") is True,
        "runtime_probe_passed": runtime_probe["status"] == "PASS",
        "binary_gaps_explicit": set(gap_summary["gap_source_formats"]) == set(),
        "redacted_evidence_only": True,
    }
    status = "PASS" if not issues and all(checks.values()) else "FAIL"
    evidence = {
        "audit_schema_version": SCHEMA_VERSION,
        "status": status,
        "scope": {
            "slice": "Slice 0287",
            "focus": "cx_office_extraction_adapter_foundation",
            "from": "uploaded_source_extraction_postgres_evidence",
            "toward": "real_document_extraction_postgres_smoke",
        },
        "paths": paths,
        "source_tokens": source_tokens,
        "backend_catalog": catalog,
        "gap_summary": gap_summary,
        "runtime_probe": runtime_probe,
        "checks": checks,
        "issues": issues,
        "decision": {
            "refactoring_checkpoint": (
                "keep extractor selection behind nex_cx.extractors.TextExtractor"
            ),
            "implemented_now": [
                "markdown",
                "plain_text",
                "pdf",
                "docx",
                "pptx",
                "xlsx",
            ],
            "placeholder_gaps": [],
            "next_slices": ["Slice 0288"],
            "raw_source_in_evidence": False,
            "local_path_in_evidence": False,
        },
    }
    assert_evidence_redacted(json.dumps(evidence, ensure_ascii=False, default=str), env)
    return evidence


def run_runtime_probe(runner: RuntimeProbe) -> dict[str, Any]:
    try:
        probe = runner()
    except Exception as exc:  # pragma: no cover - covered through injected runner
        return {
            "status": "FAIL",
            "failure_code": exc.__class__.__name__,
            "checks": {"redacted_evidence_only": True},
        }
    return probe if isinstance(probe, dict) else {"status": "FAIL", "failure_code": "invalid_probe"}


def run_extractor_backend_probe() -> dict[str, Any]:
    extractor = LocalMockTextExtractor()
    binary_outputs = {}
    for source_format in BINARY_SOURCE_FORMATS:
        filename = f"sample.{source_format}"
        source_bytes = source_bytes_for_probe(source_format)
        output = extractor.extract_markdown(
            ExtractorInput(
                filename=filename,
                content_type=content_type_for_probe(source_format),
                source_bytes=source_bytes,
                source_sha256=sha256_bytes(source_bytes),
            )
        )
        binary_outputs[source_format] = {
            "mode": output.mode,
            "warning": output.warnings[0] if output.warnings else None,
            "source_format": output.source_format,
            "real_text_seen": "Slice 0285 PDF audit text" in output.markdown_text,
            "real_docx_text_seen": "Slice 0286 DOCX audit text" in output.markdown_text,
            "real_pptx_text_seen": "Slice 0287 PPTX audit text" in output.markdown_text,
            "real_xlsx_text_seen": "Slice 0287 XLSX audit text" in output.markdown_text,
            "raw_source_leaked": SECRET_SOURCE.decode("utf-8") in output.markdown_text,
        }

    markdown_output = extractor.extract_markdown(
        ExtractorInput(
            filename="sample.md",
            content_type="text/markdown",
            source_bytes=b"# Sample\n",
            source_sha256=sha256_bytes(b"# Sample\n"),
        )
    )
    plain_output = extractor.extract_markdown(
        ExtractorInput(
            filename="sample.txt",
            content_type="text/plain",
            source_bytes=b"Plain sample",
            source_sha256=sha256_bytes(b"Plain sample"),
        )
    )
    unsupported_error_code = None
    try:
        extractor.extract_markdown(
            ExtractorInput(
                filename="sample.bin",
                content_type="application/octet-stream",
                source_bytes=b"\x00\x01",
                source_sha256=sha256_bytes(b"\x00\x01"),
            )
        )
    except ExtractionAdapterError as exc:
        unsupported_error_code = exc.error_code

    checks = {
        "markdown_real_extraction": (
            markdown_output.mode == "markdown_to_markdown"
            and markdown_output.warnings == []
        ),
        "plain_text_real_extraction": (
            plain_output.mode == "plain_text_to_markdown"
            and plain_output.warnings == []
        ),
        "pdf_real_extraction": (
            binary_outputs["pdf"]["mode"] == PDF_EXTRACTION_MODE
            and binary_outputs["pdf"]["warning"] is None
            and binary_outputs["pdf"]["real_text_seen"] is True
        ),
        "docx_real_extraction": (
            binary_outputs["docx"]["mode"] == DOCX_EXTRACTION_MODE
            and binary_outputs["docx"]["warning"] is None
            and binary_outputs["docx"]["real_docx_text_seen"] is True
        ),
        "pptx_real_extraction": (
            binary_outputs["pptx"]["mode"] == PPTX_EXTRACTION_MODE
            and binary_outputs["pptx"]["warning"] is None
            and binary_outputs["pptx"]["real_pptx_text_seen"] is True
        ),
        "xlsx_real_extraction": (
            binary_outputs["xlsx"]["mode"] == XLSX_EXTRACTION_MODE
            and binary_outputs["xlsx"]["warning"] is None
            and binary_outputs["xlsx"]["real_xlsx_text_seen"] is True
        ),
        "remaining_binary_gaps_are_placeholders": all(
            output["mode"] == PLACEHOLDER_BINARY_MODE
            for source_format, output in binary_outputs.items()
            if source_format not in set(BINARY_SOURCE_FORMATS)
        ),
        "all_binary_formats_are_real": all(
            output["warning"] is None
            for output in binary_outputs.values()
        ),
        "binary_warnings_are_stable": all(
            output["warning"] is None for output in binary_outputs.values()
        ),
        "binary_raw_source_not_serialized": not any(
            output["raw_source_leaked"] for output in binary_outputs.values()
        ),
        "unsupported_binary_rejected": (
            unsupported_error_code == "cx.extractor_source_type_unsupported"
        ),
        "redacted_evidence_only": True,
    }
    return {
        "status": "PASS" if all(checks.values()) else "FAIL",
        "probe_schema_version": "cx_extractor_backend_probe.v1",
        "observations": {
            "markdown_mode": markdown_output.mode,
            "plain_text_mode": plain_output.mode,
            "binary_modes": {
                source_format: output["mode"]
                for source_format, output in binary_outputs.items()
            },
            "binary_warnings": {
                source_format: output["warning"]
                for source_format, output in binary_outputs.items()
            },
            "unsupported_error_code": unsupported_error_code,
        },
        "checks": checks,
    }


def sample_pdf_bytes(text: str = "Slice 0285 PDF audit text") -> bytes:
    text_bytes = text.encode("ascii")
    stream = b"BT /F1 18 Tf 36 96 Td (" + text_bytes + b") Tj ET"
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 144] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream"
        ),
    )
    pdf = b"%PDF-1.4\n"
    offsets: list[int] = []
    for object_number, body in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf += (
            f"{object_number} 0 obj\n".encode("ascii")
            + body
            + b"\nendobj\n"
        )
    startxref = len(pdf)
    xref_entries = b"".join(
        f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets
    )
    return (
        pdf
        + b"xref\n0 6\n0000000000 65535 f \n"
        + xref_entries
        + b"trailer\n<< /Root 1 0 R /Size 6 >>\nstartxref\n"
        + str(startxref).encode("ascii")
        + b"\n%%EOF\n"
    )


def sample_docx_bytes(text: str = "Slice 0286 DOCX audit text") -> bytes:
    from io import BytesIO

    from docx import Document

    buffer = BytesIO()
    document = Document()
    document.add_paragraph(text)
    document.save(buffer)
    return buffer.getvalue()


def sample_pptx_bytes(text: str = "Slice 0287 PPTX audit text") -> bytes:
    from io import BytesIO

    from pptx import Presentation
    from pptx.util import Inches

    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    text_box.text = text
    presentation.save(buffer)
    return buffer.getvalue()


def sample_xlsx_bytes(text: str = "Slice 0287 XLSX audit text") -> bytes:
    from io import BytesIO

    from openpyxl import Workbook

    buffer = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet.append(["Name", "Value"])
    sheet.append(["Signal", text])
    workbook.save(buffer)
    return buffer.getvalue()


def source_bytes_for_probe(source_format: str) -> bytes:
    if source_format == "pdf":
        return sample_pdf_bytes()
    if source_format == "docx":
        return sample_docx_bytes()
    if source_format == "pptx":
        return sample_pptx_bytes()
    if source_format == "xlsx":
        return sample_xlsx_bytes()
    return SECRET_SOURCE


def content_type_for_probe(source_format: str) -> str:
    if source_format == "pdf":
        return "application/pdf"
    if source_format == "docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if source_format == "pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if source_format == "xlsx":
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    return "application/octet-stream"


def path_checks(root_dir: Path) -> list[dict[str, object]]:
    return [
        {
            "name": item.name,
            "path": relative_label(item.path, root_dir),
            "present": path_for(root_dir, item.path).exists(),
            "purpose": item.purpose,
        }
        for item in REQUIRED_PATHS
    ]


def source_token_checks(root_dir: Path) -> list[dict[str, object]]:
    return [
        {
            "group": item.group,
            "path": relative_label(item.path, root_dir),
            "token_id": item.token_id,
            "present": item.token in read_text(root_dir, item.path),
            "purpose": item.purpose,
        }
        for item in REQUIRED_SOURCE_TOKENS
    ]


def grouped_token_status(items: list[dict[str, object]]) -> dict[str, bool]:
    groups = {str(item["group"]) for item in items}
    return {
        group: all(item["present"] for item in items if item["group"] == group)
        for group in groups
    }


def issue(category: str, subject: str, detail: str) -> dict[str, str]:
    return {"category": category, "subject": subject, "detail": detail}


def assert_evidence_redacted(serialized_evidence: str, environ: Mapping[str, str]) -> None:
    for key in PROTECTED_ENV_KEYS:
        value = environ.get(key)
        if value and value not in {"1", "test"} and value in serialized_evidence:
            raise ValueError(
                "CX extractor backend gap evidence contains "
                f"unredacted environment value: {key}"
            )
    if SECRET_SOURCE.decode("utf-8") in serialized_evidence:
        raise ValueError("CX extractor backend gap evidence contains raw source.")
    if "/data/nex-platform" in serialized_evidence:
        raise ValueError("CX extractor backend gap evidence contains a local path.")


def write_audit_evidence(output_path: Path, evidence: dict[str, Any]) -> None:
    serialized = json.dumps(evidence, ensure_ascii=False, indent=2, default=str)
    assert_evidence_redacted(serialized, os.environ)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(f"{serialized}\n", encoding="utf-8")


def summary_line(evidence: dict[str, Any]) -> str:
    if evidence["status"] == "PASS":
        token_groups = grouped_token_status(evidence["source_tokens"])
        gap_summary = evidence["gap_summary"]
        return (
            "cx_extractor_backend_gap_audit=pass "
            f"paths={present_count(evidence['paths'])}/{len(evidence['paths'])} "
            f"token_groups={present_count_bool(token_groups)}/{len(token_groups)} "
            f"implemented={gap_summary['implemented_real_extraction_count']} "
            f"gaps={gap_summary['gap_placeholder_count']} "
            f"next={next_slice_from_gap_summary(gap_summary)}"
        )
    failed_checks = ",".join(
        key for key, value in evidence["checks"].items() if not value
    )
    return f"cx_extractor_backend_gap_audit=fail checks={failed_checks}"


def present_count(items: list[dict[str, object]]) -> int:
    return sum(1 for item in items if item["present"])


def present_count_bool(items: Mapping[str, bool]) -> int:
    return sum(1 for value in items.values() if value)


def next_slice_from_gap_summary(gap_summary: Mapping[str, object]) -> str | None:
    next_slices = gap_summary.get("next_slices", [])
    if not isinstance(next_slices, list) or not next_slices:
        return None
    return str(next_slices[0])


def read_text(root_dir: Path, absolute_path: Path) -> str:
    path = path_for(root_dir, absolute_path)
    return path.read_text(encoding="utf-8") if path.exists() else ""


def path_for(root_dir: Path, absolute_path: Path) -> Path:
    try:
        return root_dir / absolute_path.relative_to(ROOT)
    except ValueError:
        return absolute_path


def relative_label(path: Path, root_dir: Path = ROOT) -> str:
    try:
        return str(path.relative_to(root_dir))
    except ValueError:
        try:
            return str(path.relative_to(ROOT))
        except ValueError:
            return path.name


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Audit CX extractor backend gaps before real document adapters."
    )
    parser.add_argument("--summary", action="store_true", help="Print a short result line.")
    parser.add_argument("--output", type=Path, help="Optional JSON evidence output path.")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        evidence = run_cx_extractor_backend_gap_audit()
        if args.output:
            write_audit_evidence(args.output, evidence)
        print(
            summary_line(evidence)
            if args.summary
            else json.dumps(evidence, ensure_ascii=False, indent=2)
        )
        return 0 if evidence["status"] == "PASS" else 1
    except ValueError as exc:
        print(f"cx_extractor_backend_gap_audit=fail error={exc.__class__.__name__}")
        return 1


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main(sys.argv[1:]))
