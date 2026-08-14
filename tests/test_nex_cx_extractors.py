from __future__ import annotations

from io import BytesIO

import pytest

import pypdf
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from nex_cx.extractors import (
    BINARY_SOURCE_FORMATS,
    DOCX_EXTRACTION_MODE,
    EXTRACTED_MARKDOWN_NORMALIZATION_SCHEMA_VERSION,
    PDF_EXTRACTION_MODE,
    PPTX_EXTRACTION_MODE,
    XLSX_EXTRACTION_MODE,
    ExtractionAdapterError,
    ExtractorInput,
    ExtractorOutput,
    LocalMockTextExtractor,
    build_extracted_markdown_normalization_summary,
    classify_source_format,
    content_types_for_source_format,
    decode_utf8_source,
    docx_table_to_markdown,
    extensions_for_source_format,
    extract_docx_markdown,
    extract_pdf_markdown,
    extract_pptx_markdown,
    extract_xlsx_markdown,
    extraction_mode_for_binary_source_format,
    extractor_backend_catalog,
    extractor_backend_gap_summary,
    expected_extraction_mode_for_source_format,
    markdown_table_cell,
    markdown_from_source_text,
    mock_binary_document_markdown,
    next_slice_for_binary_source_format,
    normalize_extracted_markdown,
    normalize_extractor_output,
    rows_to_markdown_table,
    spreadsheet_cell_text,
    validate_extracted_markdown_contract,
    _ensure_trailing_newline,
)
from nex_cx.ingestion import sha256_bytes


def source(
    *,
    filename: str,
    content_type: str,
    body: bytes,
) -> ExtractorInput:
    return ExtractorInput(
        filename=filename,
        content_type=content_type,
        source_bytes=body,
        source_sha256=sha256_bytes(body),
    )


def sample_pdf_bytes(text: str = "Slice 0285 PDF extraction text") -> bytes:
    text_bytes = text.encode("ascii")
    stream = b"BT /F1 18 Tf 36 96 Td (" + text_bytes + b") Tj ET"
    objects = (
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 300 144] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        (
            b"<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream"
        ),
    )
    pdf = b"%PDF-1.4\n"
    offsets: list[int] = []
    for object_number, body in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf += (
            f"{object_number} 0 obj\n".encode("ascii")
            + body
            + b"\nendobj\n"
        )
    startxref = len(pdf)
    xref_entries = b"".join(
        f"{offset:010d} 00000 n \n".encode("ascii") for offset in offsets
    )
    return (
        pdf
        + b"xref\n0 6\n0000000000 65535 f \n"
        + xref_entries
        + b"trailer\n<< /Root 1 0 R /Size 6 >>\nstartxref\n"
        + str(startxref).encode("ascii")
        + b"\n%%EOF\n"
    )


def sample_docx_bytes(
    *,
    title: str = "Slice 0286 DOCX extraction title",
    body: str = "Slice 0286 DOCX extraction body",
) -> bytes:
    buffer = BytesIO()
    document = Document()
    document.add_paragraph(title)
    document.add_paragraph(body)
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value A"
    table.cell(1, 1).text = "Value B"
    document.save(buffer)
    return buffer.getvalue()


def sample_pptx_bytes(text: str = "Slice 0287 PPTX extraction text") -> bytes:
    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    text_box.text = text
    table = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(2),
        Inches(5),
        Inches(1),
    ).table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value A"
    table.cell(1, 1).text = "Value B"
    presentation.save(buffer)
    return buffer.getvalue()


def sample_pptx_with_ignored_empty_elements_bytes() -> bytes:
    buffer = BytesIO()
    presentation = Presentation()
    empty_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    empty_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    empty_slide.shapes.add_table(
        1,
        1,
        Inches(1),
        Inches(2),
        Inches(3),
        Inches(1),
    )
    content_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = content_slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    )
    text_box.text = "Second slide survives empty first slide"
    presentation.save(buffer)
    return buffer.getvalue()


def sample_xlsx_bytes(text: str = "Slice 0287 XLSX extraction text") -> bytes:
    buffer = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet.append(["Name", "Value"])
    sheet.append(["Signal", text])
    workbook.save(buffer)
    return buffer.getvalue()


@pytest.mark.parametrize(
    ("filename", "content_type", "expected"),
    [
        ("source.md", "application/octet-stream", "markdown"),
        ("source.txt", "application/octet-stream", "plain_text"),
        ("source.csv", "text/csv; charset=utf-8", "plain_text"),
        ("source.pdf", "application/octet-stream", "pdf"),
        ("source.docx", "application/octet-stream", "docx"),
        ("source.pptx", "application/octet-stream", "pptx"),
        ("source.xlsx", "application/octet-stream", "xlsx"),
        ("source.bin", "application/octet-stream", "unsupported"),
    ],
)
def test_classify_source_format_uses_content_type_and_extension(
    filename: str,
    content_type: str,
    expected: str,
) -> None:
    assert classify_source_format(filename=filename, content_type=content_type) == expected


def test_local_mock_extractor_preserves_markdown() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.md",
            content_type="text/markdown",
            body=b"# Existing\n\nBody",
        )
    )

    assert output.markdown_text == "# Existing\n\nBody\n"
    assert output.mode == "markdown_to_markdown"
    assert output.source_format == "markdown"
    assert output.warnings == []


def test_local_mock_extractor_wraps_plain_text() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.txt",
            content_type="text/plain",
            body=b"Plain source",
        )
    )

    assert output.markdown_text == "# source.txt\n\nPlain source\n"
    assert output.mode == "plain_text_to_markdown"
    assert output.version == "slice-0072"


def test_local_mock_extractor_extracts_pdf_text() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.pdf",
            content_type="application/pdf",
            body=sample_pdf_bytes(),
        )
    )

    assert output.markdown_text == (
        "# source.pdf\n\n## Page 1\n\nSlice 0285 PDF extraction text\n"
    )
    assert output.mode == PDF_EXTRACTION_MODE
    assert output.source_format == "pdf"
    assert output.warnings == []


def test_local_mock_extractor_extracts_docx_text_and_tables() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.docx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            body=sample_docx_bytes(),
        )
    )

    assert output.mode == DOCX_EXTRACTION_MODE
    assert output.source_format == "docx"
    assert output.warnings == []
    assert "Slice 0286 DOCX extraction title" in output.markdown_text
    assert "Slice 0286 DOCX extraction body" in output.markdown_text
    assert "## Table 1" in output.markdown_text
    assert "| Header A | Header B |" in output.markdown_text
    assert "Mock extraction placeholder." not in output.markdown_text


def test_local_mock_extractor_extracts_pptx_text_and_tables() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.pptx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            body=sample_pptx_bytes(),
        )
    )

    assert output.mode == PPTX_EXTRACTION_MODE
    assert output.source_format == "pptx"
    assert output.warnings == []
    assert "## Slide 1" in output.markdown_text
    assert "Slice 0287 PPTX extraction text" in output.markdown_text
    assert "| Header A | Header B |" in output.markdown_text
    assert "Mock extraction placeholder." not in output.markdown_text


def test_local_mock_extractor_extracts_xlsx_cells() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.xlsx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ),
            body=sample_xlsx_bytes(),
        )
    )

    assert output.mode == XLSX_EXTRACTION_MODE
    assert output.source_format == "xlsx"
    assert output.warnings == []
    assert "## Sheet Evidence" in output.markdown_text
    assert "| Name | Value |" in output.markdown_text
    assert "| Signal | Slice 0287 XLSX extraction text |" in output.markdown_text
    assert "Mock extraction placeholder." not in output.markdown_text


def test_local_mock_extractor_rejects_unsupported_source_type() -> None:
    with pytest.raises(ExtractionAdapterError) as exc:
        LocalMockTextExtractor().extract_markdown(
            source(
                filename="source.bin",
                content_type="application/octet-stream",
                body=b"\x00\x01",
            )
        )

    assert exc.value.status_code == 415
    assert exc.value.error_code == "cx.extractor_source_type_unsupported"


def test_decode_utf8_source_rejects_non_utf8_text() -> None:
    with pytest.raises(ExtractionAdapterError) as exc:
        decode_utf8_source(b"\xff\xfe")

    assert exc.value.error_code == "cx.extractor_source_encoding_unsupported"


def test_pdf_extractor_reports_parse_and_empty_text_failures() -> None:
    with pytest.raises(ExtractionAdapterError) as parse_error:
        extract_pdf_markdown(
            source(
                filename="broken.pdf",
                content_type="application/pdf",
                body=b"not a pdf",
            ),
            provider="local_mock",
            version="test",
        )
    assert parse_error.value.error_code == "cx.extractor_pdf_parse_failed"

    with pytest.raises(ExtractionAdapterError) as empty_error:
        extract_pdf_markdown(
            source(
                filename="empty.pdf",
                content_type="application/pdf",
                body=sample_pdf_bytes(""),
            ),
            provider="local_mock",
            version="test",
        )
    assert empty_error.value.error_code == "cx.extractor_pdf_text_unavailable"


def test_pdf_extractor_reports_page_text_failure(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class BrokenPage:
        def extract_text(self) -> str:
            raise RuntimeError("text extraction failed")

    class BrokenReader:
        pages = [BrokenPage()]

    monkeypatch.setattr(pypdf, "PdfReader", lambda _: BrokenReader())

    with pytest.raises(ExtractionAdapterError) as page_error:
        extract_pdf_markdown(
            source(
                filename="broken-page.pdf",
                content_type="application/pdf",
                body=sample_pdf_bytes(),
            ),
            provider="local_mock",
            version="test",
        )

    assert page_error.value.error_code == "cx.extractor_pdf_page_text_failed"
    assert page_error.value.detail.endswith("page 1.")


def test_docx_extractor_reports_parse_and_empty_text_failures() -> None:
    with pytest.raises(ExtractionAdapterError) as parse_error:
        extract_docx_markdown(
            source(
                filename="broken.docx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ),
                body=b"not a docx",
            ),
            provider="local_mock",
            version="test",
        )
    assert parse_error.value.error_code == "cx.extractor_docx_parse_failed"

    buffer = BytesIO()
    Document().save(buffer)
    with pytest.raises(ExtractionAdapterError) as empty_error:
        extract_docx_markdown(
            source(
                filename="empty.docx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                ),
                body=buffer.getvalue(),
            ),
            provider="local_mock",
            version="test",
        )
    assert empty_error.value.error_code == "cx.extractor_docx_text_unavailable"


def test_docx_extractor_ignores_empty_tables_when_paragraph_text_exists() -> None:
    buffer = BytesIO()
    document = Document()
    document.add_paragraph("Paragraph survives empty table")
    document.add_table(rows=1, cols=2)
    document.save(buffer)

    output = extract_docx_markdown(
        source(
            filename="empty-table.docx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            body=buffer.getvalue(),
        ),
        provider="local_mock",
        version="test",
    )

    assert "Paragraph survives empty table" in output.markdown_text
    assert "## Table" not in output.markdown_text


def test_pptx_extractor_reports_parse_and_empty_text_failures() -> None:
    with pytest.raises(ExtractionAdapterError) as parse_error:
        extract_pptx_markdown(
            source(
                filename="broken.pptx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                ),
                body=b"not a pptx",
            ),
            provider="local_mock",
            version="test",
        )
    assert parse_error.value.error_code == "cx.extractor_pptx_parse_failed"

    buffer = BytesIO()
    Presentation().save(buffer)
    with pytest.raises(ExtractionAdapterError) as empty_error:
        extract_pptx_markdown(
            source(
                filename="empty.pptx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                ),
                body=buffer.getvalue(),
            ),
            provider="local_mock",
            version="test",
        )
    assert empty_error.value.error_code == "cx.extractor_pptx_text_unavailable"


def test_pptx_extractor_ignores_empty_slide_elements() -> None:
    output = extract_pptx_markdown(
        source(
            filename="empty-elements.pptx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            body=sample_pptx_with_ignored_empty_elements_bytes(),
        ),
        provider="local_mock",
        version="test",
    )

    assert "## Slide 1" not in output.markdown_text
    assert "## Slide 2" in output.markdown_text
    assert "Second slide survives empty first slide" in output.markdown_text


def test_xlsx_extractor_reports_parse_and_empty_text_failures() -> None:
    with pytest.raises(ExtractionAdapterError) as parse_error:
        extract_xlsx_markdown(
            source(
                filename="broken.xlsx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                ),
                body=b"not an xlsx",
            ),
            provider="local_mock",
            version="test",
        )
    assert parse_error.value.error_code == "cx.extractor_xlsx_parse_failed"

    buffer = BytesIO()
    Workbook().save(buffer)
    with pytest.raises(ExtractionAdapterError) as empty_error:
        extract_xlsx_markdown(
            source(
                filename="empty.xlsx",
                content_type=(
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                ),
                body=buffer.getvalue(),
            ),
            provider="local_mock",
            version="test",
        )
    assert empty_error.value.error_code == "cx.extractor_xlsx_text_unavailable"


def test_markdown_helpers_are_deterministic() -> None:
    assert markdown_from_source_text("  ", filename="empty.md", content_type="text/markdown") == (
        "# empty.md\n\n"
    )
    assert markdown_from_source_text(
        "# Already\n",
        filename="already.md",
        content_type="text/markdown",
    ) == "# Already\n"
    assert mock_binary_document_markdown(
        filename="deck.pptx",
        content_type=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        source_format="pptx",
    ).startswith("# deck.pptx\n\nMock extraction placeholder.")
    assert _ensure_trailing_newline("already done\n") == "already done\n"


def test_docx_table_markdown_helpers_are_deterministic() -> None:
    assert markdown_table_cell(" A | B \n C ") == "A \\| B <br> C"
    assert spreadsheet_cell_text(None) == ""
    assert spreadsheet_cell_text(" value ") == "value"
    assert docx_table_to_markdown([["", ""], ["", ""]], table_number=1) is None
    assert docx_table_to_markdown(
        [["Name"], ["A", "B"]],
        table_number=2,
    ) == (
        "## Table 2\n\n"
        "| Name |  |\n"
        "| --- | --- |\n"
        "| A | B |"
    )
    assert rows_to_markdown_table([["H"], ["B"]], heading=None) == (
        "| H |\n| --- |\n| B |"
    )


def test_extracted_markdown_normalization_is_deterministic() -> None:
    assert normalize_extracted_markdown("# Title\r\nBody  \r\n") == "# Title\nBody\n"

    already_normalized = ExtractorOutput(
        markdown_text="# Title\n\nBody\n",
        provider="local_mock",
        mode="plain_text_to_markdown",
        version="test",
        source_format="plain_text",
        warnings=[],
    )
    assert normalize_extractor_output(already_normalized) is already_normalized

    normalized = normalize_extractor_output(
        ExtractorOutput(
            markdown_text="# Title\r\n\nBody\t\r\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        )
    )
    assert normalized.markdown_text == "# Title\n\nBody\n"


def test_extracted_markdown_normalization_summary_records_structure() -> None:
    output = LocalMockTextExtractor().extract_markdown(
        source(
            filename="source.pptx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            body=sample_pptx_bytes(),
        )
    )

    summary = build_extracted_markdown_normalization_summary(output)

    assert summary == {
        "normalization_schema_version": (
            EXTRACTED_MARKDOWN_NORMALIZATION_SCHEMA_VERSION
        ),
        "source_format": "pptx",
        "line_endings": "lf",
        "final_newline": True,
        "trailing_whitespace_present": False,
        "title_present": True,
        "required_section_heading": "slide",
        "required_section_heading_present": True,
        "heading_count": 2,
        "table_count": 1,
        "line_count": len(output.markdown_text.splitlines()),
        "char_count": len(output.markdown_text),
        "warning_count": 0,
        "contract_status": "valid",
    }


def test_extracted_markdown_contract_allows_markdown_without_h1_title() -> None:
    output = ExtractorOutput(
        markdown_text="Body only\n",
        provider="local_mock",
        mode="markdown_to_markdown",
        version="test",
        source_format="markdown",
        warnings=[],
    )

    summary = build_extracted_markdown_normalization_summary(output)

    assert summary["title_present"] is False
    assert summary["contract_status"] == "valid"


@pytest.mark.parametrize(
    "output",
    [
        ExtractorOutput(
            markdown_text="# source.pdf\n\n## Page 1\n\nBody\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="pdf",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="Body\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# source.pdf\n\nBody\n",
            provider="local_mock",
            mode=PDF_EXTRACTION_MODE,
            version="test",
            source_format="pdf",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody \n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\r\nBody\r\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[""],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody\n",
            provider="local_mock",
            mode="unknown_to_markdown",
            version="test",
            source_format="unknown",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody\n",
            provider=" ",
            mode="plain_text_to_markdown",
            version="test",
            source_format="plain_text",
            warnings=[],
        ),
        ExtractorOutput(
            markdown_text="# Title\nBody\n",
            provider="local_mock",
            mode="plain_text_to_markdown",
            version=" ",
            source_format="plain_text",
            warnings=[],
        ),
    ],
)
def test_extracted_markdown_contract_rejects_invalid_output(
    output: ExtractorOutput,
) -> None:
    with pytest.raises(ExtractionAdapterError) as exc:
        validate_extracted_markdown_contract(output)

    assert exc.value.status_code == 500
    assert exc.value.error_code == "cx.extractor_markdown_contract_invalid"


def test_extractor_backend_catalog_records_current_gaps() -> None:
    catalog = extractor_backend_catalog()
    by_format = {item.source_format: item for item in catalog}

    assert set(by_format) == {"markdown", "plain_text", *BINARY_SOURCE_FORMATS}
    assert by_format["markdown"].real_extraction is True
    assert by_format["markdown"].current_mode == "markdown_to_markdown"
    assert by_format["plain_text"].real_extraction is True
    assert by_format["plain_text"].current_mode == "plain_text_to_markdown"
    assert by_format["pdf"].real_extraction is True
    assert by_format["pdf"].current_mode == PDF_EXTRACTION_MODE
    assert by_format["pdf"].next_slice is None
    assert by_format["docx"].real_extraction is True
    assert by_format["docx"].current_mode == DOCX_EXTRACTION_MODE
    assert by_format["docx"].next_slice is None
    assert by_format["pptx"].real_extraction is True
    assert by_format["pptx"].current_mode == PPTX_EXTRACTION_MODE
    assert by_format["pptx"].next_slice is None
    assert by_format["xlsx"].real_extraction is True
    assert by_format["xlsx"].current_mode == XLSX_EXTRACTION_MODE
    assert by_format["xlsx"].next_slice is None

    summary = extractor_backend_gap_summary(catalog)
    assert summary == {
        "schema_version": "cx_extractor_backend_gap_summary.v1",
        "source_format_count": 6,
        "implemented_real_extraction_count": 6,
        "gap_placeholder_count": 0,
        "gap_source_formats": [],
        "next_slices": [],
    }


def test_extractor_backend_format_helpers_cover_text_and_unknown_formats() -> None:
    assert content_types_for_source_format("markdown") == (
        "text/markdown",
        "text/x-markdown",
    )
    assert content_types_for_source_format("plain_text") == (
        "application/json",
        "application/xml",
        "text/csv",
        "text/plain",
    )
    assert content_types_for_source_format("unknown") == ()
    assert extensions_for_source_format("markdown") == (".markdown", ".md")
    assert next_slice_for_binary_source_format("pdf") is None
    assert next_slice_for_binary_source_format("docx") is None
    assert next_slice_for_binary_source_format("pptx") is None
    assert extraction_mode_for_binary_source_format("pptx") == PPTX_EXTRACTION_MODE
    assert extraction_mode_for_binary_source_format("xlsx") == XLSX_EXTRACTION_MODE
    assert expected_extraction_mode_for_source_format("markdown") == "markdown_to_markdown"
    assert expected_extraction_mode_for_source_format("plain_text") == (
        "plain_text_to_markdown"
    )
    assert expected_extraction_mode_for_source_format("xlsx") == XLSX_EXTRACTION_MODE
    assert expected_extraction_mode_for_source_format("unknown") is None
    assert extensions_for_source_format("plain_text") == (
        ".csv",
        ".json",
        ".log",
        ".txt",
        ".xml",
    )
    assert extensions_for_source_format("unknown") == ()
