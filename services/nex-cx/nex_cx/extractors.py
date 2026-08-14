from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Protocol


MARKDOWN_CONTENT_TYPES = {"text/markdown", "text/x-markdown"}
PLAIN_TEXT_CONTENT_TYPES = {
    "application/json",
    "application/xml",
    "text/csv",
    "text/plain",
}
BINARY_DOCUMENT_CONTENT_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
}
BINARY_DOCUMENT_EXTENSIONS = {
    ".docx": "docx",
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
}
MARKDOWN_EXTENSIONS = {".markdown", ".md"}
PLAIN_TEXT_EXTENSIONS = {".csv", ".json", ".log", ".txt", ".xml"}
TEXT_SOURCE_FORMATS = ("markdown", "plain_text")
BINARY_SOURCE_FORMATS = ("pdf", "docx", "pptx", "xlsx")
SOURCE_FORMATS = (*TEXT_SOURCE_FORMATS, *BINARY_SOURCE_FORMATS)
EXTRACTED_MARKDOWN_NORMALIZATION_SCHEMA_VERSION = (
    "cx_extracted_markdown_normalization.v1"
)
PDF_EXTRACTION_MODE = "pdf_to_markdown"
DOCX_EXTRACTION_MODE = "docx_to_markdown"
PPTX_EXTRACTION_MODE = "pptx_to_markdown"
XLSX_EXTRACTION_MODE = "xlsx_to_markdown"
PLACEHOLDER_BINARY_MODE = "binary_document_placeholder_to_markdown"
PLACEHOLDER_BINARY_WARNING_PREFIX = "mock_binary_extraction_placeholder"


@dataclass(frozen=True)
class ExtractorInput:
    filename: str
    content_type: str
    source_bytes: bytes
    source_sha256: str


@dataclass(frozen=True)
class ExtractorOutput:
    markdown_text: str
    provider: str
    mode: str
    version: str
    source_format: str
    warnings: list[str]


@dataclass(frozen=True)
class ExtractorBackendCapability:
    source_format: str
    current_backend: str
    current_version: str
    current_mode: str
    status: str
    real_extraction: bool
    content_types: tuple[str, ...]
    extensions: tuple[str, ...]
    warning: str | None
    next_slice: str | None


@dataclass(frozen=True)
class ExtractionAdapterError(Exception):
    status_code: int
    error_code: str
    detail: str
    retryable: bool = False


class TextExtractor(Protocol):
    def extract_markdown(self, source: ExtractorInput) -> ExtractorOutput:
        ...


@dataclass(frozen=True)
class LocalMockTextExtractor:
    provider: str = "local_mock"
    version: str = "slice-0072"

    def extract_markdown(self, source: ExtractorInput) -> ExtractorOutput:
        source_format = classify_source_format(
            filename=source.filename,
            content_type=source.content_type,
        )
        if source_format in TEXT_SOURCE_FORMATS:
            source_text = decode_utf8_source(source.source_bytes)
            return ExtractorOutput(
                markdown_text=markdown_from_source_text(
                    source_text,
                    filename=source.filename,
                    content_type=source.content_type,
                ),
                provider=self.provider,
                mode=f"{source_format}_to_markdown",
                version=self.version,
                source_format=source_format,
                warnings=[],
            )
        if source_format == "pdf":
            return extract_pdf_markdown(source, provider=self.provider, version=self.version)
        if source_format == "docx":
            return extract_docx_markdown(source, provider=self.provider, version=self.version)
        if source_format == "pptx":
            return extract_pptx_markdown(source, provider=self.provider, version=self.version)
        if source_format == "xlsx":
            return extract_xlsx_markdown(source, provider=self.provider, version=self.version)
        raise ExtractionAdapterError(
            status_code=415,
            error_code="cx.extractor_source_type_unsupported",
            detail=f"No extractor is registered for source type: {source.content_type}",
        )


def classify_source_format(*, filename: str, content_type: str) -> str:
    normalized_content_type = content_type.split(";", maxsplit=1)[0].strip().lower()
    suffix = Path(filename).suffix.lower()
    if normalized_content_type in MARKDOWN_CONTENT_TYPES or suffix in MARKDOWN_EXTENSIONS:
        return "markdown"
    if (
        normalized_content_type in PLAIN_TEXT_CONTENT_TYPES
        or normalized_content_type.startswith("text/")
        or suffix in PLAIN_TEXT_EXTENSIONS
    ):
        return "plain_text"
    if normalized_content_type in BINARY_DOCUMENT_CONTENT_TYPES:
        return BINARY_DOCUMENT_CONTENT_TYPES[normalized_content_type]
    if suffix in BINARY_DOCUMENT_EXTENSIONS:
        return BINARY_DOCUMENT_EXTENSIONS[suffix]
    return "unsupported"


def extractor_backend_catalog(
    *,
    provider: str = "local_mock",
    version: str = "slice-0072",
) -> tuple[ExtractorBackendCapability, ...]:
    return (
        ExtractorBackendCapability(
            source_format="markdown",
            current_backend=provider,
            current_version=version,
            current_mode="markdown_to_markdown",
            status="implemented",
            real_extraction=True,
            content_types=tuple(sorted(MARKDOWN_CONTENT_TYPES)),
            extensions=tuple(sorted(MARKDOWN_EXTENSIONS)),
            warning=None,
            next_slice=None,
        ),
        ExtractorBackendCapability(
            source_format="plain_text",
            current_backend=provider,
            current_version=version,
            current_mode="plain_text_to_markdown",
            status="implemented",
            real_extraction=True,
            content_types=tuple(sorted(PLAIN_TEXT_CONTENT_TYPES)),
            extensions=tuple(sorted(PLAIN_TEXT_EXTENSIONS)),
            warning=None,
            next_slice=None,
        ),
        *(
            binary_extractor_capability(
                source_format,
                provider=provider,
                version=version,
            )
            for source_format in BINARY_SOURCE_FORMATS
        ),
    )


def extractor_backend_gap_summary(
    catalog: tuple[ExtractorBackendCapability, ...] | None = None,
) -> dict[str, object]:
    capabilities = catalog if catalog is not None else extractor_backend_catalog()
    implemented = [item for item in capabilities if item.real_extraction]
    gaps = [item for item in capabilities if not item.real_extraction]
    return {
        "schema_version": "cx_extractor_backend_gap_summary.v1",
        "source_format_count": len(capabilities),
        "implemented_real_extraction_count": len(implemented),
        "gap_placeholder_count": len(gaps),
        "gap_source_formats": [item.source_format for item in gaps],
        "next_slices": [
            item.next_slice for item in gaps if item.next_slice is not None
        ],
    }


def content_types_for_source_format(source_format: str) -> tuple[str, ...]:
    if source_format == "markdown":
        return tuple(sorted(MARKDOWN_CONTENT_TYPES))
    if source_format == "plain_text":
        return tuple(sorted(PLAIN_TEXT_CONTENT_TYPES))
    return tuple(
        sorted(
            content_type
            for content_type, mapped_format in BINARY_DOCUMENT_CONTENT_TYPES.items()
            if mapped_format == source_format
        )
    )


def extensions_for_source_format(source_format: str) -> tuple[str, ...]:
    if source_format == "markdown":
        return tuple(sorted(MARKDOWN_EXTENSIONS))
    if source_format == "plain_text":
        return tuple(sorted(PLAIN_TEXT_EXTENSIONS))
    return tuple(
        sorted(
            extension
            for extension, mapped_format in BINARY_DOCUMENT_EXTENSIONS.items()
            if mapped_format == source_format
        )
    )


def next_slice_for_binary_source_format(source_format: str) -> str | None:
    return None


def binary_extractor_capability(
    source_format: str,
    *,
    provider: str,
    version: str,
) -> ExtractorBackendCapability:
    return ExtractorBackendCapability(
        source_format=source_format,
        current_backend=provider,
        current_version=version,
        current_mode=extraction_mode_for_binary_source_format(source_format),
        status="implemented",
        real_extraction=True,
        content_types=content_types_for_source_format(source_format),
        extensions=extensions_for_source_format(source_format),
        warning=None,
        next_slice=None,
    )


def extraction_mode_for_binary_source_format(source_format: str) -> str:
    return {
        "pdf": PDF_EXTRACTION_MODE,
        "docx": DOCX_EXTRACTION_MODE,
        "pptx": PPTX_EXTRACTION_MODE,
        "xlsx": XLSX_EXTRACTION_MODE,
    }[source_format]


def expected_extraction_mode_for_source_format(source_format: str) -> str | None:
    if source_format == "markdown":
        return "markdown_to_markdown"
    if source_format == "plain_text":
        return "plain_text_to_markdown"
    if source_format in BINARY_SOURCE_FORMATS:
        return extraction_mode_for_binary_source_format(source_format)
    return None


def normalize_extracted_markdown(markdown_text: str) -> str:
    normalized = markdown_text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = "\n".join(line.rstrip(" \t") for line in normalized.split("\n"))
    return _ensure_trailing_newline(normalized)


def normalize_extractor_output(output: ExtractorOutput) -> ExtractorOutput:
    normalized_markdown = normalize_extracted_markdown(output.markdown_text)
    if normalized_markdown == output.markdown_text:
        return output
    return ExtractorOutput(
        markdown_text=normalized_markdown,
        provider=output.provider,
        mode=output.mode,
        version=output.version,
        source_format=output.source_format,
        warnings=list(output.warnings),
    )


def build_extracted_markdown_normalization_summary(
    output: ExtractorOutput,
) -> dict[str, object]:
    lines = output.markdown_text.splitlines()
    title_present = output.markdown_text.startswith("# ")
    heading_count = sum(1 for line in lines if is_markdown_heading_line(line))
    table_count = count_markdown_tables(lines)
    required_heading = required_section_heading_for_source_format(output.source_format)
    required_heading_present = (
        any(line.startswith(required_heading) for line in lines)
        if required_heading is not None
        else False
    )
    summary: dict[str, object] = {
        "normalization_schema_version": (
            EXTRACTED_MARKDOWN_NORMALIZATION_SCHEMA_VERSION
        ),
        "source_format": output.source_format,
        "line_endings": "lf" if "\r" not in output.markdown_text else "mixed",
        "final_newline": output.markdown_text.endswith("\n"),
        "trailing_whitespace_present": has_trailing_whitespace(lines),
        "title_present": title_present,
        "required_section_heading": required_heading_name(output.source_format),
        "required_section_heading_present": required_heading_present,
        "heading_count": heading_count,
        "table_count": table_count,
        "line_count": len(lines),
        "char_count": len(output.markdown_text),
        "warning_count": len(output.warnings),
        "contract_status": "valid",
    }
    validate_extracted_markdown_contract(output, summary)
    return summary


def validate_extracted_markdown_contract(
    output: ExtractorOutput,
    summary: dict[str, object] | None = None,
) -> None:
    computed = (
        summary
        if summary is not None
        else build_extracted_markdown_normalization_summary(output)
    )
    if output.source_format not in SOURCE_FORMATS:
        raise_extracted_markdown_contract_error(
            f"unsupported source_format: {output.source_format}"
        )
    expected_mode = expected_extraction_mode_for_source_format(output.source_format)
    if expected_mode is not None and output.mode != expected_mode:
        raise_extracted_markdown_contract_error(
            f"mode {output.mode} does not match source_format {output.source_format}"
        )
    if not str(output.provider).strip():
        raise_extracted_markdown_contract_error("provider must be a non-empty string")
    if not str(output.version).strip():
        raise_extracted_markdown_contract_error("version must be a non-empty string")
    if not output.markdown_text.endswith("\n"):
        raise_extracted_markdown_contract_error("final newline is required")
    if "\r" in output.markdown_text:
        raise_extracted_markdown_contract_error("line endings must be LF")
    if computed["trailing_whitespace_present"] is True:
        raise_extracted_markdown_contract_error(
            "trailing spaces and tabs are not allowed"
        )
    if output.source_format in ("plain_text", *BINARY_SOURCE_FORMATS) and (
        computed["title_present"] is not True
    ):
        raise_extracted_markdown_contract_error(
            "plain text and binary extraction output must start with an H1 title"
        )
    required_heading = required_section_heading_for_source_format(output.source_format)
    if required_heading is not None and (
        computed["required_section_heading_present"] is not True
    ):
        raise_extracted_markdown_contract_error(
            f"{output.source_format} output must include {required_heading.strip()}"
        )
    if any(not warning.strip() for warning in output.warnings):
        raise_extracted_markdown_contract_error(
            "warnings must be non-empty strings when present"
        )


def required_section_heading_for_source_format(source_format: str) -> str | None:
    return {
        "pdf": "## Page ",
        "pptx": "## Slide ",
        "xlsx": "## Sheet ",
    }.get(source_format)


def required_heading_name(source_format: str) -> str:
    return {
        "pdf": "page",
        "pptx": "slide",
        "xlsx": "sheet",
    }.get(source_format, "none")


def is_markdown_heading_line(line: str) -> bool:
    marker_length = len(line) - len(line.lstrip("#"))
    return 1 <= marker_length <= 6 and line[marker_length : marker_length + 1] == " "


def count_markdown_tables(lines: list[str]) -> int:
    return sum(
        1
        for index, line in enumerate(lines[:-1])
        if line.startswith("| ")
        and line.endswith(" |")
        and lines[index + 1].startswith("| ---")
    )


def has_trailing_whitespace(lines: list[str]) -> bool:
    return any(line.endswith((" ", "\t")) for line in lines)


def raise_extracted_markdown_contract_error(detail: str) -> None:
    raise ExtractionAdapterError(
        status_code=500,
        error_code="cx.extractor_markdown_contract_invalid",
        detail=f"Extracted Markdown contract violation: {detail}.",
    )


def decode_utf8_source(source_bytes: bytes) -> str:
    try:
        return source_bytes.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ExtractionAdapterError(
            status_code=415,
            error_code="cx.extractor_source_encoding_unsupported",
            detail="Text source bytes must be UTF-8 encoded.",
        ) from exc


def markdown_from_source_text(
    source_text: str,
    *,
    filename: str,
    content_type: str,
) -> str:
    stripped = source_text.strip()
    if not stripped:
        return f"# {filename}\n\n"
    if filename.lower().endswith(".md") or content_type == "text/markdown":
        return _ensure_trailing_newline(stripped)
    return _ensure_trailing_newline(f"# {filename}\n\n{stripped}")


def mock_binary_document_markdown(
    *,
    filename: str,
    content_type: str,
    source_format: str,
) -> str:
    return (
        f"# {filename}\n\n"
        "Mock extraction placeholder.\n\n"
        f"- source_format: {source_format}\n"
        f"- content_type: {content_type}\n"
    )


def extract_pdf_markdown(
    source: ExtractorInput,
    *,
    provider: str,
    version: str,
) -> ExtractorOutput:
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(source.source_bytes))
    except Exception as exc:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_pdf_parse_failed",
            detail="PDF source could not be parsed.",
        ) from exc

    page_sections: list[str] = []
    warnings: list[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            raise ExtractionAdapterError(
                status_code=422,
                error_code="cx.extractor_pdf_page_text_failed",
                detail=f"PDF page text extraction failed: page {page_number}.",
            ) from exc
        normalized = page_text.strip()
        if not normalized:
            warnings.append(f"pdf_page_text_empty:{page_number}")
            continue
        page_sections.append(f"## Page {page_number}\n\n{normalized}")
    if not page_sections:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_pdf_text_unavailable",
            detail="PDF source did not contain extractable text.",
        )
    return ExtractorOutput(
        markdown_text=_ensure_trailing_newline(
            f"# {source.filename}\n\n" + "\n\n".join(page_sections)
        ),
        provider=provider,
        mode=PDF_EXTRACTION_MODE,
        version=version,
        source_format="pdf",
        warnings=warnings,
    )


def extract_docx_markdown(
    source: ExtractorInput,
    *,
    provider: str,
    version: str,
) -> ExtractorOutput:
    try:
        from docx import Document

        document = Document(BytesIO(source.source_bytes))
    except Exception as exc:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_docx_parse_failed",
            detail="DOCX source could not be parsed.",
        ) from exc

    sections: list[str] = []
    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]
    if paragraphs:
        sections.append("\n\n".join(paragraphs))
    for table_number, table in enumerate(document.tables, start=1):
        table_markdown = docx_table_to_markdown(
            [
                [cell.text.strip() for cell in row.cells]
                for row in table.rows
            ],
            table_number=table_number,
        )
        if table_markdown is not None:
            sections.append(table_markdown)
    if not sections:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_docx_text_unavailable",
            detail="DOCX source did not contain extractable text.",
        )
    return ExtractorOutput(
        markdown_text=_ensure_trailing_newline(
            f"# {source.filename}\n\n" + "\n\n".join(sections)
        ),
        provider=provider,
        mode=DOCX_EXTRACTION_MODE,
        version=version,
        source_format="docx",
        warnings=[],
    )


def extract_pptx_markdown(
    source: ExtractorInput,
    *,
    provider: str,
    version: str,
) -> ExtractorOutput:
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(source.source_bytes))
    except Exception as exc:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_pptx_parse_failed",
            detail="PPTX source could not be parsed.",
        ) from exc

    sections: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                table_markdown = rows_to_markdown_table(
                    [
                        [cell.text.strip() for cell in row.cells]
                        for row in shape.table.rows
                    ],
                    heading=None,
                )
                if table_markdown is not None:
                    slide_parts.append(table_markdown)
        if slide_parts:
            sections.append(f"## Slide {slide_number}\n\n" + "\n\n".join(slide_parts))
    if not sections:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_pptx_text_unavailable",
            detail="PPTX source did not contain extractable text.",
        )
    return ExtractorOutput(
        markdown_text=_ensure_trailing_newline(
            f"# {source.filename}\n\n" + "\n\n".join(sections)
        ),
        provider=provider,
        mode=PPTX_EXTRACTION_MODE,
        version=version,
        source_format="pptx",
        warnings=[],
    )


def extract_xlsx_markdown(
    source: ExtractorInput,
    *,
    provider: str,
    version: str,
) -> ExtractorOutput:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(
            BytesIO(source.source_bytes),
            data_only=True,
            read_only=True,
        )
    except Exception as exc:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_xlsx_parse_failed",
            detail="XLSX source could not be parsed.",
        ) from exc

    sections: list[str] = []
    for sheet in workbook.worksheets:
        table_markdown = rows_to_markdown_table(
            [
                [spreadsheet_cell_text(cell.value) for cell in row]
                for row in sheet.iter_rows()
            ],
            heading=f"## Sheet {sheet.title}",
        )
        if table_markdown is not None:
            sections.append(table_markdown)
    if not sections:
        raise ExtractionAdapterError(
            status_code=422,
            error_code="cx.extractor_xlsx_text_unavailable",
            detail="XLSX source did not contain extractable text.",
        )
    return ExtractorOutput(
        markdown_text=_ensure_trailing_newline(
            f"# {source.filename}\n\n" + "\n\n".join(sections)
        ),
        provider=provider,
        mode=XLSX_EXTRACTION_MODE,
        version=version,
        source_format="xlsx",
        warnings=[],
    )


def docx_table_to_markdown(
    rows: list[list[str]],
    *,
    table_number: int,
) -> str | None:
    return rows_to_markdown_table(rows, heading=f"## Table {table_number}")


def rows_to_markdown_table(
    rows: list[list[str]],
    *,
    heading: str | None,
) -> str | None:
    non_empty_rows = [row for row in rows if any(cell.strip() for cell in row)]
    if not non_empty_rows:
        return None
    width = max(len(row) for row in non_empty_rows)
    padded_rows = [
        [markdown_table_cell(cell) for cell in [*row, *[""] * (width - len(row))]]
        for row in non_empty_rows
    ]
    header = padded_rows[0]
    body = padded_rows[1:]
    table_lines = [] if heading is None else [heading, ""]
    table_lines.extend(
        [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]
    )
    table_lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(table_lines)


def markdown_table_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", "<br>").strip()


def spreadsheet_cell_text(value: object) -> str:
    return "" if value is None else str(value).strip()


def _ensure_trailing_newline(value: str) -> str:
    if value.endswith("\n"):
        return value
    return f"{value}\n"
